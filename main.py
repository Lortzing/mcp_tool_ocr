# document_parser_refactored.py
from __future__ import annotations

import io
import os
import time
import json
import base64
import traceback
import urllib.parse
import asyncio
from dataclasses import dataclass
from typing import Optional, Dict, Any, List

import fitz
import pdfplumber
import camelot
import pytesseract
from pytesseract import Output as TESSERACT_OUTPUT
from PIL import Image
import docx
import pandas as pd
import openpyxl
from pptx import Presentation
import httpx
from openai import OpenAI

from config import (
    OPENAI_API_KEY,
    OPENAI_BASE_URL,
    QWEN_VLM_MODEL,
)

from fastmcp import FastMCP
from starlette.requests import Request
from starlette.responses import JSONResponse

def _b64(bytestr: bytes) -> str:
    return base64.b64encode(bytestr).decode("utf-8")


def _binary_placeholder(label: str) -> str:
    """Generate a consistent placeholder for removed binary payloads."""

    clean_label = label.strip() or "content"
    return f"[binary content omitted: {clean_label}]"


class VLMClient:
    """Minimal wrapper around the OpenAI client for multimodal interactions."""

    def __init__(
        self,
        api_key: Optional[str] = None,
        base_url: Optional[str] = None,
        model: Optional[str] = None,
        timeout: float = 120.0,
    ):
        self.api_key = api_key or ""
        self.base_url = base_url or ""
        self.model = model or "qwen3-vl-plus"
        self.timeout = timeout
        self._client: OpenAI | None = None
        if self.api_key:
            self._client = OpenAI(api_key=self.api_key, base_url=self.base_url, timeout=timeout)

    def _ensure_client(self) -> bool:
        if self._client is None:
            return False
        return True

    def _safe_dump(self, response: Any) -> Any:
        if response is None:
            return None
        dump = getattr(response, "model_dump", None)
        if callable(dump):
            try:
                return dump()
            except Exception:
                pass
        json_fn = getattr(response, "json", None)
        if callable(json_fn):
            try:
                return json.loads(json_fn())
            except Exception:
                pass
        return getattr(response, "__dict__", str(response))

    def _extract_text(self, response: Any) -> str:
        if response is None:
            return ""
        output = getattr(response, "output", None)
        if isinstance(output, list):
            collected: List[str] = []
            for item in output:
                contents = getattr(item, "content", None)
                if isinstance(contents, list):
                    for content in contents:
                        text = getattr(content, "text", None) or getattr(content, "value", None)
                        if isinstance(text, str):
                            collected.append(text)
            if collected:
                return "".join(collected)

        choices = getattr(response, "choices", None)
        if isinstance(choices, list):
            texts = []
            for choice in choices:
                message = getattr(choice, "message", None)
                if isinstance(message, dict):
                    content = message.get("content")
                    if isinstance(content, str):
                        texts.append(content)
                elif hasattr(message, "content"):
                    content = getattr(message, "content")
                    if isinstance(content, str):
                        texts.append(content)
            if texts:
                return "".join(texts)
        return ""

    def generate(self, prompt: str, images_b64: List[str]) -> Dict[str, Any]:
        if not self._ensure_client():
            return {"error": "OpenAI client not configured"}

        content: List[Dict[str, Any]] = [{"type": "input_text", "text": prompt}]
        for img in images_b64:
            content.append({"type": "input_image", "image": {"base64": img}})

        try:
            response = self._client.responses.create(
                model=self.model,
                input=[{"role": "user", "content": content}],
            )
            return {
                "markdown": self._extract_text(response),
                "raw": self._safe_dump(response),
            }
        except Exception as exc:
            return {"error": f"OpenAI call failed: {exc}"}


@dataclass
class VLMRequest:
    prompt: str
    images_b64: List[str]


class VLMTaskManager:
    """Manage asynchronous VLM invocations with a concurrency ceiling."""

    def __init__(self, client: Optional[VLMClient], max_concurrent: int = 5):
        self.client = client
        self.max_concurrent = max(1, min(max_concurrent, 5))

    async def _run_all(self, requests: List[VLMRequest]) -> List[Dict[str, Any]]:
        semaphore = asyncio.Semaphore(self.max_concurrent)

        async def _task(request: VLMRequest) -> Dict[str, Any]:
            if not self.client:
                return {"error": "OpenAI client not configured"}
            if not request.images_b64:
                return {"error": "No images provided for VLM request"}
            async with semaphore:
                try:
                    return await asyncio.to_thread(
                        self.client.generate, request.prompt, request.images_b64
                    )
                except Exception as exc:
                    return {"error": f"OpenAI call failed: {exc}"}

        tasks = [_task(req) for req in requests]
        return await asyncio.gather(*tasks)

    def run_requests(self, requests: List[VLMRequest]) -> List[Dict[str, Any]]:
        if not requests:
            return []
        if not self.client:
            return [{"error": "OpenAI client not configured"} for _ in requests]
        return asyncio.run(self._run_all(requests))

    def run_single(self, request: VLMRequest) -> Dict[str, Any]:
        results = self.run_requests([request])
        return results[0] if results else {"error": "OpenAI client not configured"}


class OCRHelper:
    @staticmethod
    def _normalize_lang(lang: str | List[str]) -> str:
        """
        规范化语言参数：
        - 如果传入 list，例如 ["eng", "chi_sim"]，则转成 "eng+chi_sim"
        - 如果传入 str，原样返回
        """
        if isinstance(lang, (list, tuple)):
            return "+".join(lang)
        return lang

    @staticmethod
    def image_to_data(image_bytes: bytes, lang: str | List[str] = "eng") -> Dict[str, Any]:
        """
        OCR 识别图像，支持多语言。
        e.g. lang="eng+chi_sim" 或 ["eng", "chi_sim"]
        """
        lang_str = OCRHelper._normalize_lang(lang)
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        try:
            data: dict = pytesseract.image_to_data(img, lang=lang_str)
        except pytesseract.TesseractError as e:
            return {"error": f"OCR failed with lang={lang_str}: {str(e)}"}

        out = []
        n = len(data.get("text", []))
        for i in range(n):
            txt = str(data.get("text", [])[i] or "").strip()
            if not txt:
                continue
            conf_raw = data.get("conf", [])[i]
            try:
                conf = float(conf_raw) if conf_raw != "-1" else None
            except Exception:
                conf = None
            bbox = {
                "left": int(data.get("left", [])[i]),
                "top": int(data.get("top", [])[i]),
                "width": int(data.get("width", [])[i]),
                "height": int(data.get("height", [])[i]),
            }
            out.append({"text": txt, "conf": conf, "bbox": bbox})
        return {"blocks": out, "lang": lang_str}



# -------------------------
# Parsers (clean interface: parse(file_bytes: bytes))
# -------------------------

class BaseParser:
    def __init__(
        self,
        use_ocr: bool = False,
        vlm_client: Optional[VLMClient] = None,
        vlm_manager: Optional[VLMTaskManager] = None,
        ocr_lang: str | List[str] = "eng",
    ):
        # OCR is disabled by default. We keep the flag for backwards compatibility
        # but avoid invoking OCR logic in new flows unless explicitly requested.
        self.use_ocr = use_ocr
        self.vlm_client = vlm_client
        self.vlm_manager = vlm_manager or (VLMTaskManager(vlm_client) if vlm_client else None)
        self.ocr_lang = ocr_lang

    def parse(self, data: bytes) -> Dict[str, Any]:
        raise NotImplementedError


class PDFParser(BaseParser):
    def __init__(
        self,
        *args,
        camelot_enable: bool = True,
        pdfplumber_enable: bool = True,
        parse_mode: str = "ocr",
        **kwargs,
    ):
        super().__init__(*args, **kwargs)
        self.camelot_enable = camelot_enable
        self.pdfplumber_enable = pdfplumber_enable
        self.parse_mode = parse_mode.lower()
        if self.parse_mode not in {"ocr", "vlm"}:
            self.parse_mode = "ocr"

    def _extract_selectable_text(self, doc) -> List[Dict[str, Any]]:
        pages = []
        for i, page in enumerate(doc):
            try:
                text = page.get_text("text", sort=True)
            except Exception:
                text = ""
            normalized = text or ""
            elements: List[Dict[str, Any]] = []
            if normalized.strip():
                elements.append({"type": "text", "order": 1, "text": normalized})
            pages.append({"page_number": i + 1, "text": normalized, "elements": elements})
        return pages

    def _extract_tables_pdfplumber(self, pdf_bytes: bytes) -> Dict[str, Any]:
        out = []
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                try:
                    tables = page.extract_tables()
                except Exception:
                    tables = []
                parsed = [t for t in tables]
                out.append({"page_number": i + 1, "tables": parsed})
        return {"pages": out}

    def _extract_tables_camelot(self, pdf_bytes: bytes) -> Dict[str, Any]:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
            tf.write(pdf_bytes)
            tmpname = tf.name
        try:
            tables = camelot.read_pdf(tmpname, pages="all")
            extracted = []
            for t in tables:
                try:
                    df = t.df
                    extracted.append({"shape": df.shape, "csv": df.to_csv(index=False)})
                except Exception:
                    try:
                        extracted.append({"csv": t.to_csv()})
                    except Exception:
                        extracted.append({"error": "failed to convert camelot table"})
            return {"tables_count": len(extracted), "tables": extracted}
        finally:
            try:
                os.remove(tmpname)
            except Exception:
                pass

    def _render_page_to_png(self, page, zoom: float = 2.0) -> bytes:
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        return pix.tobytes("png")

    def _extract_embedded_images(self, doc) -> List[Dict[str, Any]]:
        images_out = []
        for pno in range(len(doc)):
            page = doc[pno]
            try:
                imgs = page.get_images(full=True)
            except Exception:
                imgs = []
            for img_ref in imgs:
                try:
                    xref = img_ref[0]
                    fitz.Pixmap(doc, xref)  # ensure the reference is valid
                    images_out.append(
                        {
                            "page": pno + 1,
                            "description": _binary_placeholder("embedded image"),
                        }
                    )
                except Exception:
                    continue
        return images_out

    @staticmethod
    def _append_page_element(page_entry: Dict[str, Any], element: Dict[str, Any]) -> None:
        elements = page_entry.setdefault("elements", [])
        element["order"] = len(elements) + 1
        elements.append(element)

    def parse(self, data: bytes) -> Dict[str, Any]:
        start = time.time()
        try:
            doc = fitz.open(stream=data, filetype="pdf")
        except Exception as e:
            return {"error": "failed to open pdf with pymupdf", "details": str(e)}

        if self.parse_mode == "vlm":
            return self._parse_with_vlm(doc, start)
        return self._parse_with_ocr(data, doc, start)

    def _parse_with_ocr(self, pdf_bytes: bytes, doc, start: float) -> Dict[str, Any]:
        result: Dict[str, Any] = {
            "code": 0,
            "type": "pdf",
            "mode": "ocr",
            "pages": [],
            "tables": {},
            "images": [],
            "warnings": [],
        }

        try:
            result["pages"] = self._extract_selectable_text(doc)
        except Exception as e:
            result["warnings"].append(f"pymupdf text extraction failed: {str(e)}")

        if self.pdfplumber_enable:
            try:
                pdfplumber_tables = self._extract_tables_pdfplumber(pdf_bytes)
                result["tables_pdfplumber"] = pdfplumber_tables
                for page_tables in pdfplumber_tables.get("pages", []):
                    page_number = page_tables.get("page_number")
                    if not isinstance(page_number, int):
                        continue
                    index = page_number - 1
                    if 0 <= index < len(result.get("pages", [])):
                        page_entry = result["pages"][index]
                        for table in page_tables.get("tables", []):
                            self._append_page_element(
                                page_entry,
                                {"type": "table", "source": "pdfplumber", "table": table},
                            )
            except Exception as e:
                result["warnings"].append(f"pdfplumber failed: {str(e)}")

        if self.camelot_enable:
            try:
                camelot_tables = self._extract_tables_camelot(pdf_bytes)
                result["tables_camelot"] = camelot_tables
            except Exception as e:
                result["warnings"].append(f"camelot failed: {str(e)}")

        need_ocr_pages: List[int] = []
        try:
            pages = result.get("pages", [])
            if not pages:
                need_ocr_pages = list(range(len(doc)))
            else:
                need_ocr_pages = [i for i, page in enumerate(pages) if not page.get("text", "").strip()]
        except Exception:
            need_ocr_pages = list(range(len(doc)))

        vlm_page_results: List[Dict[str, Any]] = []
        vlm_requests: List[VLMRequest] = []
        vlm_request_metadata: List[Dict[str, Any]] = []
        for pnum in need_ocr_pages:
            try:
                page = doc[pnum]
                img_bytes = self._render_page_to_png(page, zoom=2.0)
                if self.use_ocr:
                    try:
                        ocr_result = OCRHelper.image_to_data(img_bytes, lang=self.ocr_lang)
                    except Exception as ocr_err:
                        ocr_result = {"error": str(ocr_err)}
                    while len(result["pages"]) <= pnum:
                        result["pages"].append({"page_number": len(result["pages"]) + 1, "text": "", "elements": []})
                    result["pages"][pnum]["ocr"] = ocr_result
                    self._append_page_element(result["pages"][pnum], {"type": "ocr_text", "ocr": ocr_result})

                while len(result["pages"]) <= pnum:
                    result["pages"].append({"page_number": len(result["pages"]) + 1, "text": "", "elements": []})

                page_entry = result["pages"][pnum]

                if self.vlm_manager:
                    prompt = self._build_pdf_page_prompt(pnum + 1)
                    placeholder = {"type": "vlm_markdown", "vlm": None}
                    self._append_page_element(page_entry, placeholder)
                    vlm_requests.append(VLMRequest(prompt=prompt, images_b64=[_b64(img_bytes)]))
                    vlm_request_metadata.append({"page_number": pnum + 1, "placeholder": placeholder})
                else:
                    vlm_response = {"error": "OpenAI client not configured"}
                    self._append_page_element(page_entry, {"type": "vlm_markdown", "vlm": vlm_response})
                    vlm_page_results.append({"page_number": pnum + 1, "vlm": vlm_response})
            except Exception as e:
                result["warnings"].append(f"analysis failed on page {pnum + 1}: {str(e)}")

        try:
            result["images"] = self._extract_embedded_images(doc)
        except Exception as e:
            result["warnings"].append(f"extract embedded images failed: {str(e)}")
        else:
            for img in result.get("images", []):
                page_number = img.get("page")
                if not isinstance(page_number, int):
                    continue
                index = page_number - 1
                if 0 <= index < len(result.get("pages", [])):
                    page_entry = result["pages"][index]
                    self._append_page_element(
                        page_entry,
                        {
                            "type": "image",
                            "description": _binary_placeholder("embedded image"),
                        },
                    )

        if vlm_requests and self.vlm_manager:
            responses = self.vlm_manager.run_requests(vlm_requests)
            for meta, resp in zip(vlm_request_metadata, responses):
                meta_placeholder = meta.get("placeholder")
                if isinstance(meta_placeholder, dict):
                    meta_placeholder["vlm"] = resp
                vlm_page_results.append({"page_number": meta.get("page_number"), "vlm": resp})

        if vlm_page_results:
            result["vlm_pages"] = vlm_page_results

        result["elapsed_seconds"] = time.time() - start
        return result

    def _build_pdf_page_prompt(self, page_number: int) -> str:
        return (
            "You are converting scanned PDF pages into structured Markdown. "
            f"Transcribe all readable content on page {page_number} from the provided image. "
            "Use Markdown headings, lists and tables whenever appropriate, and render mathematics "
            "using LaTeX syntax between $...$ or $$...$$. Describe figures or charts in one or two "
            "sentences placed where they appear in the page. Respond with Markdown only."
        )

    def _render_document_to_images(self, doc) -> List[Dict[str, Any]]:
        images: List[Dict[str, Any]] = []
        for idx in range(len(doc)):
            try:
                png_bytes = self._render_page_to_png(doc[idx], zoom=2.0)
                images.append(
                    {
                        "page_number": idx + 1,
                        "image_b64": _b64(png_bytes),
                        "placeholder": _binary_placeholder(f"page {idx + 1} render"),
                    }
                )
            except Exception:
                images.append({"page_number": idx + 1, "error": "render_failed"})
        return images

    def _build_pdf_document_prompt(self, page_images: List[Dict[str, Any]]) -> str:
        total_pages = len(page_images)
        return (
            "You are an expert assistant that converts PDF documents into high quality Markdown. "
            f"The user will provide {total_pages} page image(s). Process them in order, producing "
            "Markdown that mirrors the layout and reading order. Follow these rules:\n"
            "1. Preserve all textual content faithfully.\n"
            "2. Use Markdown tables for tabular data.\n"
            "3. Replace charts, figures or photos with concise descriptive captions at the same position.\n"
            "4. Render mathematics using LaTeX syntax inside $...$ or $$...$$.\n"
            "5. Start each page with a heading like '## Page X' to maintain pagination.\n"
            "Return only Markdown without additional commentary."
        )

    def _parse_with_vlm(self, doc, start: float) -> Dict[str, Any]:
        page_images_raw = self._render_document_to_images(doc)
        images_payload = [item.get("image_b64") for item in page_images_raw if item.get("image_b64")]

        if not images_payload:
            return {
                "error": "pdf to image conversion failed",
                "details": page_images_raw,
            }

        pages_metadata = []
        page_images: List[Dict[str, Any]] = []
        for item in page_images_raw:
            page_number = item.get("page_number")
            entry = {"page_number": page_number, "elements": []}
            placeholder = item.get("placeholder")
            if not placeholder and item.get("image_b64"):
                placeholder = _binary_placeholder(f"page {page_number} render")
            if placeholder:
                self._append_page_element(
                    entry,
                    {"type": "image", "description": placeholder},
                )
                page_images.append({"page_number": page_number, "description": placeholder})
            elif item.get("error"):
                page_images.append({"page_number": page_number, "error": item.get("error")})
            pages_metadata.append(entry)

        if not self.vlm_manager:
            vlm_response = {"error": "OpenAI client not configured"}
        else:
            prompt = self._build_pdf_document_prompt(page_images_raw)
            vlm_response = self.vlm_manager.run_single(VLMRequest(prompt=prompt, images_b64=images_payload))

        document_elements = []
        if vlm_response:
            document_elements.append({"type": "vlm_markdown", "order": 1, "vlm": vlm_response})

        return {
            "code": 0,
            "type": "pdf",
            "mode": "vlm",
            "pages": pages_metadata,
            "page_images": page_images,
            "vlm": vlm_response,
            "document_elements": document_elements,
            "elapsed_seconds": time.time() - start,
        }


class ImageParser(BaseParser):
    def parse(self, data: bytes) -> Dict[str, Any]:
        start = time.time()
        result: Dict[str, Any] = {"code": 0, "type": "image", "vlm": None, "elapsed_seconds": None}
        img_b64 = _b64(data)
        placeholder = _binary_placeholder("original image")
        content_flow: List[Dict[str, Any]] = [
            {"type": "image", "order": 1, "description": placeholder}
        ]
        result["image"] = {"description": placeholder}

        if self.use_ocr:
            try:
                result["ocr"] = OCRHelper.image_to_data(data, lang=self.ocr_lang)
            except Exception as e:
                result["ocr"] = {"error": str(e)}
            content_flow.append({"type": "ocr_text", "order": len(content_flow) + 1, "ocr": result["ocr"]})

        if self.vlm_manager:
            try:
                prompt = self._build_image_prompt()
                vlm_res = self.vlm_manager.run_single(VLMRequest(prompt=prompt, images_b64=[img_b64]))
                result["vlm"] = vlm_res
            except Exception as e:
                result["vlm"] = {"error": str(e)}
        else:
            result["vlm"] = {"error": "OpenAI client not configured"}
        content_flow.append({"type": "vlm_markdown", "order": len(content_flow) + 1, "vlm": result["vlm"]})
        result["content_flow"] = content_flow
        result["elapsed_seconds"] = time.time() - start
        return result

    @staticmethod
    def _build_image_prompt() -> str:
        return (
            "You are analysing a single image. Extract every readable text segment, "
            "describe the visual elements and layout, and reproduce any tables using "
            "Markdown table syntax. When mathematical expressions appear, rewrite them "
            "using LaTeX math enclosed in $...$ or $$...$$. Respond with concise Markdown "
            "combining the transcription and the scene description."
        )


class DocxParser(BaseParser):
    @staticmethod
    def _iter_blocks(doc) -> List[Any]:
        try:
            from docx.oxml.table import CT_Tbl  # type: ignore
            from docx.oxml.text.paragraph import CT_P  # type: ignore
            from docx.table import Table as DocxTable  # type: ignore
            from docx.text.paragraph import Paragraph  # type: ignore

            if hasattr(doc, "element") and hasattr(doc.element, "body"):
                for child in doc.element.body.iterchildren():  # type: ignore[attr-defined]
                    if isinstance(child, CT_P):
                        yield "paragraph", Paragraph(child, doc)
                    elif isinstance(child, CT_Tbl):
                        yield "table", DocxTable(child, doc)
                return
        except Exception:
            pass

        for para in getattr(doc, "paragraphs", []):
            yield "paragraph", para
        for table in getattr(doc, "tables", []):
            yield "table", table

    def parse(self, data: bytes) -> Dict[str, Any]:
        start = time.time()
        try:
            doc = docx.Document(io.BytesIO(data))
            paragraphs: List[str] = []
            tables: List[List[List[str]]] = []
            flow: List[Dict[str, Any]] = []

            for block_type, block in self._iter_blocks(doc):
                if block_type == "paragraph":
                    text = getattr(block, "text", "")
                    cleaned = text.strip()
                    if cleaned:
                        idx = len(paragraphs)
                        paragraphs.append(cleaned)
                        flow.append({"type": "paragraph", "index": idx, "order": len(flow) + 1, "text": cleaned})
                elif block_type == "table":
                    rows: List[List[str]] = []
                    for r in getattr(block, "rows", []):
                        row_cells = [getattr(c, "text", "") for c in getattr(r, "cells", [])]
                        rows.append(row_cells)
                    idx = len(tables)
                    tables.append(rows)
                    flow.append({"type": "table", "index": idx, "order": len(flow) + 1, "rows": rows})

            images: List[Dict[str, Any]] = []
            images_vlm: List[Dict[str, Any]] = []
            vlm_jobs: List[Dict[str, Any]] = []
            rels_iterable: List[Any] = []
            try:
                part = getattr(doc, "part", None)
                rels = getattr(part, "rels", None)
                if rels is not None:
                    rels_iterable = list(rels.values())  # type: ignore[call-arg]
            except Exception:
                rels_iterable = []

            for rel in rels_iterable:
                try:
                    if "image" in getattr(rel, "reltype", ""):
                        blob = rel.target_part.blob
                        image_b64 = _b64(blob)
                        placeholder = _binary_placeholder("embedded image")
                        image_entry = {
                            "index": len(images),
                            "description": placeholder,
                            "vlm": None,
                        }
                        images.append(image_entry)
                        flow.append(
                            {
                                "type": "image",
                                "index": image_entry["index"],
                                "order": len(flow) + 1,
                                "description": placeholder,
                            }
                        )
                        if self.vlm_manager:
                            prompt = ImageParser._build_image_prompt()
                            vlm_jobs.append(
                                {
                                    "request": VLMRequest(prompt=prompt, images_b64=[image_b64]),
                                    "image": image_entry,
                                    "order": len(flow),
                                }
                            )
                        else:
                            error_res = {"error": "OpenAI client not configured"}
                            image_entry["vlm"] = error_res
                            images_vlm.append({"order": len(flow), "vlm": error_res})
                except Exception as exc:
                    images_vlm.append({"error": f"failed to analyse images: {str(exc)}"})

            if vlm_jobs and self.vlm_manager:
                responses = self.vlm_manager.run_requests([job["request"] for job in vlm_jobs])
                for job, resp in zip(vlm_jobs, responses):
                    job["image"]["vlm"] = resp
                    images_vlm.append({"order": job["order"], "vlm": resp})

            if images_vlm:
                images_vlm.sort(key=lambda item: item.get("order", float("inf")))

            return {
                "code": 0,
                "type": "docx",
                "content": {"paragraphs": paragraphs, "tables": tables, "flow": flow},
                "images": images,
                "images_vlm": images_vlm,
                "elapsed_seconds": time.time() - start,
            }
        except Exception as e:
            return {"error": str(e)}


class XlsxParser(BaseParser):
    def parse(self, data: bytes) -> Dict[str, Any]:
        start = time.time()
        try:
            sheets = pd.read_excel(io.BytesIO(data), sheet_name=None)
            out = {}
            for name, df in sheets.items():
                out[name] = {"shape": df.shape, "csv": df.to_csv(index=False)}
            return {"code": 0, "type": "xlsx", "content": {"sheets": out}, "elapsed_seconds": time.time() - start}
        except Exception as e:
            return {"error": str(e)}


class PptxParser(BaseParser):
    """Parse .pptx files using python-pptx"""
    def parse(self, data: bytes) -> Dict[str, Any]:
        start = time.time()
        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as e:
            return {"error": f"failed to open pptx: {str(e)}"}

        slides_out = []
        images_out = []
        image_jobs: List[Dict[str, Any]] = []
        for i, slide in enumerate(prs.slides):
            slide_text_chunks = []
            notes_text = None
            elements: List[Dict[str, Any]] = []
            try:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_text_chunks.append(text)
                            elements.append({"type": "text", "order": len(elements) + 1, "text": text})

                    # 提取图片
                    try:
                        image = getattr(shape, "image", None)
                        if image is not None:
                            blob = image.blob
                            img_b64 = _b64(blob)
                            placeholder = _binary_placeholder("slide image")
                            image_entry = {
                                "slide": i + 1,
                                "description": placeholder,
                                "vlm": None,
                            }
                            images_out.append(image_entry)
                            image_element = {
                                "type": "image",
                                "order": len(elements) + 1,
                                "description": placeholder,
                                "image_index": len(images_out) - 1,
                            }
                            elements.append(image_element)
                            if self.vlm_manager:
                                prompt = ImageParser._build_image_prompt()
                                image_jobs.append(
                                    {
                                        "request": VLMRequest(prompt=prompt, images_b64=[img_b64]),
                                        "entry": image_entry,
                                        "element": image_element,
                                        "slide": i + 1,
                                    }
                                )
                            else:
                                error_res = {"error": "OpenAI client not configured"}
                                image_entry["vlm"] = error_res
                                image_element["vlm"] = error_res
                    except Exception:
                        pass

                try:
                    if getattr(slide, "notes_slide", None) is not None and getattr(slide.notes_slide, "notes_text_frame", None) is not None:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text:
                            elements.append({"type": "notes", "order": len(elements) + 1, "text": notes_text})
                except Exception:
                    notes_text = None
            except Exception:
                pass
            slides_out.append({"slide_number": i + 1, "texts": slide_text_chunks, "notes": notes_text, "elements": elements})

        images_vlm_results: List[Dict[str, Any]] = []
        if image_jobs and self.vlm_manager:
            responses = self.vlm_manager.run_requests([job["request"] for job in image_jobs])
            for job, resp in zip(image_jobs, responses):
                job["entry"]["vlm"] = resp
                job["element"]["vlm"] = resp
                images_vlm_results.append({"slide": job["slide"], "vlm": resp})
        else:
            for img in images_out:
                vlm_res = img.get("vlm") or {"error": "OpenAI client not configured"}
                images_vlm_results.append({"slide": img.get("slide"), "vlm": vlm_res})

        if images_vlm_results:
            images_vlm_results.sort(key=lambda item: item.get("slide", 0))

        return {
            "code": 0,
            "type": "pptx",
            "slides": slides_out,
            "images": images_out,
            "images_vlm": images_vlm_results,
            "elapsed_seconds": time.time() - start,
        }


# -------------------------
# DocumentProcessor: picks parser by extension and dispatches
# -------------------------

class DocumentProcessor:
    def __init__(
        self,
        use_ocr: bool = False,
        ocr_lang: str = "eng",
        camelot_enable: bool = True,
        pdfplumber_enable: bool = True,
        pdf_parse_mode: str = "ocr",
    ):
        resolved_base_url = OPENAI_BASE_URL
        resolved_token = OPENAI_API_KEY
        resolved_model = QWEN_VLM_MODEL

        self.vlm_client = VLMClient(
            api_key=resolved_token,
            base_url=resolved_base_url,
            model=resolved_model,
        )
        self.vlm_manager = VLMTaskManager(self.vlm_client)

        self.use_ocr = use_ocr
        self.ocr_lang = ocr_lang
        self.camelot_enable = camelot_enable
        self.pdfplumber_enable = pdfplumber_enable
        self.pdf_parse_mode = pdf_parse_mode.lower() if isinstance(pdf_parse_mode, str) else "ocr"
        if self.pdf_parse_mode not in {"ocr", "vlm"}:
            self.pdf_parse_mode = "ocr"

    def _select_parser(self, filename: str) -> BaseParser:
        lower = filename.lower()
        if lower.endswith(".pdf"):
            return PDFParser(
                use_ocr=self.use_ocr,
                vlm_client=self.vlm_client,
                vlm_manager=self.vlm_manager,
                ocr_lang=self.ocr_lang,
                camelot_enable=self.camelot_enable,
                pdfplumber_enable=self.pdfplumber_enable,
                parse_mode=self.pdf_parse_mode,
            )
        if any(lower.endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]):
            return ImageParser(
                use_ocr=self.use_ocr,
                vlm_client=self.vlm_client,
                vlm_manager=self.vlm_manager,
                ocr_lang=self.ocr_lang,
            )
        if lower.endswith(".docx") or lower.endswith(".doc"):
            return DocxParser(
                use_ocr=self.use_ocr,
                vlm_client=self.vlm_client,
                vlm_manager=self.vlm_manager,
                ocr_lang=self.ocr_lang,
            )
        if lower.endswith(".xlsx") or lower.endswith(".xls"):
            return XlsxParser(
                use_ocr=self.use_ocr,
                vlm_client=self.vlm_client,
                vlm_manager=self.vlm_manager,
                ocr_lang=self.ocr_lang,
            )
        if lower.endswith(".pptx") or lower.endswith(".ppt"):
            return PptxParser(
                use_ocr=self.use_ocr,
                vlm_client=self.vlm_client,
                vlm_manager=self.vlm_manager,
                ocr_lang=self.ocr_lang,
            )
        raise ValueError("Unsupported file extension")

    def parse(self, file_bytes: bytes, filename: str) -> Dict[str, Any]:
        parser = self._select_parser(filename)
        return parser.parse(file_bytes)


# -------------------------
# Markdown formatting utilities
# -------------------------

def _format_vlm_markdown(vlm_payload: Any) -> str:
    if isinstance(vlm_payload, dict):
        markdown = vlm_payload.get("markdown")
        if isinstance(markdown, str) and markdown.strip():
            return markdown.strip()
        error = vlm_payload.get("error")
        if isinstance(error, str):
            return f"Error: {error}"
    if isinstance(vlm_payload, str) and vlm_payload.strip():
        return vlm_payload.strip()
    return ""


def _format_ocr_blocks(ocr_payload: Any, limit: int = 8) -> List[str]:
    if not isinstance(ocr_payload, dict):
        return []
    blocks = ocr_payload.get("blocks") or []
    if not isinstance(blocks, list):
        return []
    lines: List[str] = []
    lang = ocr_payload.get("lang")
    if isinstance(lang, str) and lang:
        lines.append(f"- OCR language: {lang}")
    count = 0
    for block in blocks:
        if count >= limit:
            break
        text = ""
        if isinstance(block, dict):
            text = str(block.get("text", "")).strip()
        elif isinstance(block, str):
            text = block.strip()
        if text:
            lines.append(f"- {text}")
            count += 1
    total = len(blocks)
    if total > count:
        lines.append(f"- ... ({total} OCR segments in total)")
    return lines


def _format_table_rows(rows: Any) -> str:
    if not isinstance(rows, list) or not rows:
        return ""
    max_cols = max((len(row) for row in rows if isinstance(row, list)), default=0)
    if max_cols == 0:
        return ""
    header = rows[0] if isinstance(rows[0], list) else [rows[0]]
    header = list(header) + [""] * (max_cols - len(header))
    body_rows = []
    for row in rows[1:]:
        if isinstance(row, list):
            padded = list(row) + [""] * (max_cols - len(row))
            body_rows.append(padded)
        else:
            body_rows.append([str(row)] + [""] * (max_cols - 1))
    divider = "|".join([" --- " for _ in range(max_cols)])
    lines = [
        "| " + " | ".join(str(cell).strip() for cell in header) + " |",
        "|" + divider + "|",
    ]
    for row in body_rows:
        lines.append("| " + " | ".join(str(cell).strip() for cell in row) + " |")
    return "\n".join(lines)


def _format_pdf_markdown(parsed: Dict[str, Any]) -> List[str]:
    lines: List[str] = []
    mode = parsed.get("mode")
    pages = parsed.get("pages") or []
    lines.append("## PDF Overview")
    if isinstance(mode, str):
        lines.append(f"- Mode: {mode}")
    lines.append(f"- Pages parsed: {len(pages)}")
    lines.append("")

    warnings = parsed.get("warnings") or []
    if warnings:
        lines.append("### Warnings")
        for warning in warnings:
            lines.append(f"- {warning}")
        lines.append("")

    for page in pages:
        page_number = page.get("page_number") if isinstance(page, dict) else None
        if page_number is None:
            continue
        lines.append(f"### Page {page_number}")
        text_content = ""
        if isinstance(page, dict):
            text_content = str(page.get("text", "")).strip()
        if text_content:
            lines.append("#### Extracted Text")
            lines.append(text_content)
            lines.append("")

        ocr_lines = _format_ocr_blocks(page.get("ocr")) if isinstance(page, dict) else []
        if ocr_lines:
            lines.append("#### OCR Highlights")
            lines.extend(ocr_lines)
            lines.append("")

        elements = page.get("elements") if isinstance(page, dict) else []
        if elements:
            for element in elements:
                if not isinstance(element, dict):
                    continue
                etype = element.get("type")
                if etype == "table":
                    source = element.get("source", "table")
                    lines.append(f"- Table ({source})")
                    md_table = _format_table_rows(element.get("table"))
                    if md_table:
                        lines.append(md_table)
                    lines.append("")
                elif etype == "image":
                    lines.append(f"- {element.get('description', _binary_placeholder('image'))}")
                elif etype == "vlm_markdown":
                    vlm_text = _format_vlm_markdown(element.get("vlm"))
                    if vlm_text:
                        lines.append("#### VLM Markdown")
                        lines.append(vlm_text)
                        lines.append("")
        lines.append("")

    tables_pdfplumber = parsed.get("tables_pdfplumber") or {}
    if isinstance(tables_pdfplumber, dict) and tables_pdfplumber.get("pages"):
        lines.append("### Tables (pdfplumber)")
        for table_page in tables_pdfplumber.get("pages", []):
            if not isinstance(table_page, dict):
                continue
            page_number = table_page.get("page_number")
            tables = table_page.get("tables") or []
            lines.append(f"- Page {page_number}: {len(tables)} table(s)")
        lines.append("")

    camelot_tables = parsed.get("tables_camelot") or {}
    if isinstance(camelot_tables, dict) and camelot_tables.get("tables"):
        lines.append("### Tables (camelot)")
        tables = camelot_tables.get("tables", [])
        lines.append(f"- Extracted tables: {len(tables)}")
        lines.append("")

    images = parsed.get("images") or []
    if images:
        lines.append("### Embedded Images")
        for image in images:
            if not isinstance(image, dict):
                continue
            page_number = image.get("page")
            description = image.get("description", _binary_placeholder("embedded image"))
            if page_number is not None:
                lines.append(f"- Page {page_number}: {description}")
        lines.append("")

    vlm_pages = parsed.get("vlm_pages") or []
    if vlm_pages:
        lines.append("### VLM Page Summaries")
        for entry in vlm_pages:
            if not isinstance(entry, dict):
                continue
            page_number = entry.get("page_number")
            vlm_text = _format_vlm_markdown(entry.get("vlm"))
            if vlm_text:
                lines.append(f"#### Page {page_number}")
                lines.append(vlm_text)
        lines.append("")

    vlm_document = _format_vlm_markdown(parsed.get("vlm"))
    if vlm_document:
        lines.append("### Document-Level Markdown")
        lines.append(vlm_document)
        lines.append("")

    return lines


def _format_image_markdown(parsed: Dict[str, Any]) -> List[str]:
    lines = ["## Image Overview"]
    placeholder = parsed.get("image", {}).get("description")
    if placeholder:
        lines.append(f"- Image: {placeholder}")
    lines.append("")
    if parsed.get("ocr"):
        ocr_lines = _format_ocr_blocks(parsed.get("ocr"))
        if ocr_lines:
            lines.append("### OCR Highlights")
            lines.extend(ocr_lines)
            lines.append("")
    vlm_text = _format_vlm_markdown(parsed.get("vlm"))
    if vlm_text:
        lines.append("### VLM Markdown")
        lines.append(vlm_text)
        lines.append("")
    return lines


def _format_docx_markdown(parsed: Dict[str, Any]) -> List[str]:
    lines = ["## DOCX Overview"]
    content = parsed.get("content") or {}
    paragraphs = content.get("paragraphs") if isinstance(content, dict) else []
    if paragraphs:
        lines.append("### Paragraphs")
        for idx, para in enumerate(paragraphs, 1):
            lines.append(f"{idx}. {para}")
        lines.append("")
    tables = content.get("tables") if isinstance(content, dict) else []
    if tables:
        lines.append("### Tables")
        for idx, table in enumerate(tables, 1):
            lines.append(f"#### Table {idx}")
            md_table = _format_table_rows(table)
            if md_table:
                lines.append(md_table)
            lines.append("")
    images = parsed.get("images") or []
    if images:
        lines.append("### Embedded Images")
        for idx, image in enumerate(images, 1):
            description = image.get("description", _binary_placeholder("embedded image")) if isinstance(image, dict) else _binary_placeholder("embedded image")
            lines.append(f"- Image {idx}: {description}")
        lines.append("")
    images_vlm = parsed.get("images_vlm") or []
    if images_vlm:
        lines.append("### Image VLM Summaries")
        for entry in images_vlm:
            if not isinstance(entry, dict):
                continue
            order = entry.get("order")
            vlm_text = _format_vlm_markdown(entry.get("vlm"))
            if vlm_text:
                lines.append(f"- Flow position {order}: {vlm_text}")
        lines.append("")
    return lines


def _format_xlsx_markdown(parsed: Dict[str, Any]) -> List[str]:
    lines = ["## XLSX Overview"]
    content = parsed.get("content") or {}
    sheets = content.get("sheets") if isinstance(content, dict) else {}
    if sheets:
        for name, sheet in sheets.items():
            lines.append(f"### Sheet: {name}")
            if isinstance(sheet, dict):
                shape = sheet.get("shape")
                if shape:
                    lines.append(f"- Shape: {shape}")
                csv_text = sheet.get("csv")
                if isinstance(csv_text, str) and csv_text.strip():
                    lines.append("```csv")
                    lines.append(csv_text.strip())
                    lines.append("```")
            lines.append("")
    return lines


def _format_pptx_markdown(parsed: Dict[str, Any]) -> List[str]:
    lines = ["## PPTX Overview"]
    slides = parsed.get("slides") or []
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_number = slide.get("slide_number")
        lines.append(f"### Slide {slide_number}")
        texts = slide.get("texts") or []
        if texts:
            lines.append("#### Text Content")
            for idx, text in enumerate(texts, 1):
                lines.append(f"{idx}. {text}")
            lines.append("")
        notes = slide.get("notes")
        if notes:
            lines.append("#### Speaker Notes")
            lines.append(notes)
            lines.append("")
        elements = slide.get("elements") or []
        for element in elements:
            if not isinstance(element, dict):
                continue
            if element.get("type") == "image":
                description = element.get("description", _binary_placeholder("slide image"))
                lines.append(f"- {description}")
            elif element.get("type") == "vlm_markdown":
                vlm_text = _format_vlm_markdown(element.get("vlm"))
                if vlm_text:
                    lines.append("#### Image VLM Markdown")
                    lines.append(vlm_text)
                    lines.append("")
            elif element.get("type") == "notes":
                continue
        lines.append("")
    images_vlm = parsed.get("images_vlm") or []
    if images_vlm:
        lines.append("### Image VLM Summaries")
        for entry in images_vlm:
            if not isinstance(entry, dict):
                continue
            slide_number = entry.get("slide")
            vlm_text = _format_vlm_markdown(entry.get("vlm"))
            if vlm_text:
                lines.append(f"- Slide {slide_number}: {vlm_text}")
        lines.append("")
    return lines


def build_markdown_report(parsed: Any, filename: Optional[str] = None) -> str:
    if isinstance(parsed, dict) and parsed.get("error"):
        lines = ["# Document Parse Error"]
        if filename:
            lines.append(f"- **File:** {filename}")
        lines.append(f"- **Error:** {parsed.get('error')}")
        details = parsed.get("details") or parsed.get("trace")
        if isinstance(details, str) and details.strip():
            lines.append("")
            lines.append("```")
            lines.append(details.strip())
            lines.append("```")
        return "\n".join(lines)

    if not isinstance(parsed, dict):
        return str(parsed)

    doc_type = str(parsed.get("type", "document")).upper()
    title = filename or "document"
    lines: List[str] = [f"# Parse Result for {title}", "", f"- **Type:** {doc_type}"]
    elapsed = parsed.get("elapsed_seconds")
    if isinstance(elapsed, (int, float)):
        lines.append(f"- **Elapsed Seconds:** {elapsed:.2f}")
    lines.append("")

    type_dispatch = {
        "PDF": _format_pdf_markdown,
        "IMAGE": _format_image_markdown,
        "DOCX": _format_docx_markdown,
        "XLSX": _format_xlsx_markdown,
        "PPTX": _format_pptx_markdown,
    }
    formatter = type_dispatch.get(doc_type)
    if formatter:
        lines.extend(formatter(parsed))
    else:
        lines.append("## Content")
        lines.append(json.dumps(parsed, ensure_ascii=False, indent=2))

    return "\n".join(line for line in lines if line is not None)


# -------------------------
# MCP wrapper
# -------------------------

mcp = FastMCP() if FastMCP else None

if mcp:
    @mcp.tool
    def parse_document_url(
        file_url: str,
        options: Optional[Dict[str, Any]] = None,
    ) -> Any:
        """Main MCP entrypoint."""

        opts = options or {}
        try:
            with httpx.Client(timeout=opts.get("timeout", 60.0)) as client:
                resp = client.get(file_url)
                resp.raise_for_status()
                file_bytes = resp.content
        except Exception as e:
            return {"error": f"Failed to fetch file: {str(e)}"}

        filename = os.path.basename(urllib.parse.urlparse(file_url).path) or "downloaded_file"
        
        dp = DocumentProcessor(
            use_ocr=opts.get("use_ocr", False),
            ocr_lang=opts.get("ocr_lang", "eng"),
            camelot_enable=opts.get("camelot_enable", True),
            pdfplumber_enable=opts.get("pdfplumber_enable", True),
            pdf_parse_mode=opts.get("pdf_parse_mode", "ocr"),
        )

        try:
            parsed = dp.parse(file_bytes, filename)
            return build_markdown_report(parsed, filename=filename)
        except ValueError as e:
            return {"error": str(e)}
        except Exception:
            return {"error": "Unhandled exception during parse", "trace": traceback.format_exc()}

    @mcp.custom_route("/api/health", methods=["GET"])
    async def health(request: Request) -> JSONResponse:
        return JSONResponse({"status": "ok", "service": "document-parser-mcp-refactored"})

if __name__ == "__main__":
    if mcp:
        host = os.getenv("MCP_HOST", "localhost")
        port = int(os.getenv("MCP_PORT", "50002"))
        mcp.run(transport="http", host=host, port=port, path="/mcp", log_level="info")
    else:
        print("FastMCP not available, cannot run service.")
